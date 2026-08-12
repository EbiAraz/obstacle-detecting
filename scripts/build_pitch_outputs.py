from pathlib import Path
import argparse
import re
from datetime import datetime
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas
from reportlab.lib.units import cm

try:
    from pptx import Presentation
    from pptx.util import Pt
except Exception:
    Presentation = None


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description='Build HTML/PDF/PPTX pitch outputs from markdown.')
    parser.add_argument(
        '--input',
        default='PITCH_DECK_RAILSAFE_TURKEY.md',
        help='Input markdown file path',
    )
    parser.add_argument(
        '--basename',
        default='PITCH_DECK_RAILSAFE_TURKEY',
        help='Base name for generated outputs (without extension)',
    )
    parser.add_argument(
        '--footer',
        default='Ankara-Istanbul Corridor Pitch Deck',
        help='Footer text used in PPTX slides',
    )
    parser.add_argument(
        '--lang',
        default='en',
        help='Language code for deck localization (e.g., en, tr)',
    )
    parser.add_argument(
        '--badge',
        default='RailSafe AI',
        help='Top badge text on slides',
    )
    parser.add_argument(
        '--slide-label',
        default='',
        help='Optional slide number label (e.g., "Slide", "Slayt"). Empty means number only.',
    )
    parser.add_argument(
        '--generated-label',
        default='',
        help='Optional generated timestamp label. Empty uses localized default.',
    )
    return parser.parse_args()


def parse_slides(text: str):
    lines = text.splitlines()
    slides = []
    current_title = 'RailSafe AI'
    current_lines = []

    for line in lines:
        if line.startswith('## '):
            if current_lines:
                slides.append((current_title, current_lines))
            current_title = line[3:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines:
        slides.append((current_title, current_lines))

    return slides


def localized_defaults(lang: str, generated_label_override: str, slide_label_override: str):
    lang = (lang or 'en').lower()
    if lang.startswith('tr'):
        generated_label = generated_label_override or 'Olusturulma'
        slide_label = slide_label_override or 'Slayt'
        timestamp = datetime.now().strftime('%d.%m.%Y %H:%M')
    else:
        generated_label = generated_label_override or 'Generated'
        slide_label = slide_label_override or 'Slide'
        timestamp = datetime.now().strftime('%Y-%m-%d %H:%M')
    return generated_label, slide_label, timestamp


def md_to_html_line(line: str) -> str:
    line = line.rstrip()
    if not line:
        return ''
    if line.startswith('### '):
        return f"<h3>{line[4:].strip()}</h3>"
    if line.startswith('- '):
        return f"<li>{line[2:].strip()}</li>"
    if re.match(r'^\d+\.\s+', line):
        content = re.sub(r'^\d+\.\s+', '', line)
        return f"<li>{content}</li>"
    return f"<p>{line}</p>"

def build_html(slides, html_path: Path, *, lang: str, badge_text: str, slide_label: str, generated_line: str):
    slide_html = []
    for i, (title, body_lines) in enumerate(slides, start=1):
        body_html = []
        in_ul = False
        in_ol = False

        for ln in body_lines:
            if ln.startswith('- '):
                if in_ol:
                    body_html.append('</ol>')
                    in_ol = False
                if not in_ul:
                    body_html.append('<ul>')
                    in_ul = True
                body_html.append(md_to_html_line(ln))
            elif re.match(r'^\d+\.\s+', ln):
                if in_ul:
                    body_html.append('</ul>')
                    in_ul = False
                if not in_ol:
                    body_html.append('<ol>')
                    in_ol = True
                body_html.append(md_to_html_line(ln))
            else:
                if in_ul:
                    body_html.append('</ul>')
                    in_ul = False
                if in_ol:
                    body_html.append('</ol>')
                    in_ol = False
                converted = md_to_html_line(ln)
                if converted:
                    body_html.append(converted)

        if in_ul:
            body_html.append('</ul>')
        if in_ol:
            body_html.append('</ol>')

        idx_text = f"{slide_label} {i:02d}" if slide_label else f"{i:02d}"
        slide_html.append(f"""
<section class=\"slide\">\
  <div class=\"slide-head\">\
    <span class=\"badge\">{badge_text}</span>\
    <span class=\"idx\">{idx_text}</span>\
  </div>\
  <h2>{title}</h2>\
  <div class=\"content\">{''.join(body_html)}</div>\
  <div class=\"footer-note\">{generated_line}</div>\
</section>
""")

    html = f"""<!doctype html>
<html lang=\"{lang}\">
<head>
<meta charset=\"utf-8\" />
<meta name=\"viewport\" content=\"width=device-width, initial-scale=1\" />
<title>RailSafe AI Pitch Deck</title>
<style>
:root {{
  --bg: #0b1020;
  --panel: #131a2f;
  --text: #e8edf8;
  --muted: #a8b4d1;
  --accent: #22c55e;
  --accent2: #38bdf8;
}}
* {{ box-sizing: border-box; }}
body {{
  margin: 0;
  font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
  background: radial-gradient(circle at 20% 10%, #172140 0%, var(--bg) 45%, #070b16 100%);
  color: var(--text);
}}
.deck {{
  width: min(1200px, 96vw);
  margin: 18px auto 40px;
}}
.slide {{
  min-height: 88vh;
  margin: 0 0 18px;
  background: linear-gradient(165deg, rgba(19,26,47,.95), rgba(15,20,36,.98));
  border: 1px solid #283556;
  border-radius: 22px;
  padding: 36px 44px;
  box-shadow: 0 20px 50px rgba(0,0,0,.35);
  page-break-after: always;
}}
.slide-head {{ display: flex; justify-content: space-between; align-items: center; margin-bottom: 20px; }}
.badge {{ font-weight: 700; color: #06121b; background: linear-gradient(90deg, var(--accent), var(--accent2)); padding: 8px 14px; border-radius: 999px; font-size: .9rem; }}
.idx {{ color: var(--muted); font-size: .95rem; letter-spacing: .08em; }}
h2 {{ font-size: clamp(1.7rem, 3vw, 2.5rem); margin: 0 0 14px; line-height: 1.2; }}
h3 {{ color: #8bd9ff; margin: 12px 0 6px; }}
p, li {{ font-size: clamp(1.02rem, 1.5vw, 1.2rem); line-height: 1.55; color: var(--text); }}
ul, ol {{ margin-top: 10px; }}
.content p {{ color: #d8e1f4; }}
.footer-note {{
    margin-top: 18px;
    color: #8a9bc0;
    font-size: .88rem;
    letter-spacing: .02em;
}}
@media print {{
  body {{ background: #fff; }}
  .slide {{ min-height: auto; border-radius: 0; box-shadow: none; margin: 0; border: none; }}
}}
</style>
</head>
<body>
<div class=\"deck\">{''.join(slide_html)}</div>
</body>
</html>
"""
    html_path.write_text(html, encoding='utf-8')


def build_pdf(slides, pdf_path: Path, *, slide_label: str, generated_line: str):
    c = canvas.Canvas(str(pdf_path), pagesize=A4)
    _w, h = A4
    margin_x = 1.8 * cm
    margin_y = 1.8 * cm

    for idx, (title, body_lines) in enumerate(slides, start=1):
        y = h - margin_y
        c.setFont('Helvetica-Bold', 22)
        prefix = f"{slide_label} {idx:02d}" if slide_label else f"{idx:02d}"
        c.drawString(margin_x, y, f"{prefix}. {title}")
        y -= 1.1 * cm

        c.setFont('Helvetica', 11)
        for raw in body_lines:
            line = raw.strip()
            if not line:
                y -= 0.25 * cm
                continue

            if line.startswith('### '):
                txt = line[4:].strip()
                c.setFont('Helvetica-Bold', 13)
                if y < margin_y + 2 * cm:
                    c.showPage()
                    y = h - margin_y
                c.drawString(margin_x, y, txt)
                y -= 0.6 * cm
                c.setFont('Helvetica', 11)
                continue

            if line.startswith('- '):
                txt = '- ' + line[2:].strip()
            elif re.match(r'^\d+\.\s+', line):
                txt = re.sub(r'^\d+\.\s+', '- ', line)
            else:
                txt = line

            chunks = []
            max_chars = 105
            while len(txt) > max_chars:
                cut = txt.rfind(' ', 0, max_chars)
                if cut == -1:
                    cut = max_chars
                chunks.append(txt[:cut])
                txt = txt[cut:].lstrip()
            chunks.append(txt)

            for chunk in chunks:
                if y < margin_y + 1.5 * cm:
                    c.showPage()
                    y = h - margin_y
                    c.setFont('Helvetica', 11)
                c.drawString(margin_x, y, chunk)
                y -= 0.5 * cm

        c.setFont('Helvetica', 9)
        c.drawString(margin_x, margin_y - 0.2 * cm, generated_line)
        c.showPage()

    c.save()


def build_pptx(slides, pptx_path: Path, footer_text: str, *, badge_text: str, slide_label: str, generated_line: str):
    if Presentation is None:
        return False

    import pptx

    prs = Presentation()
    prs.slide_width = 12192000   # 13.333 in
    prs.slide_height = 6858000   # 7.5 in

    for idx, (title, body_lines) in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = pptx.dml.color.RGBColor(9, 20, 40)

        top = slide.shapes.add_textbox(Pt(24), Pt(14), Pt(880), Pt(26))
        tf_top = top.text_frame
        tf_top.clear()
        p_top = tf_top.paragraphs[0]
        r_top = p_top.add_run()
        r_top.text = badge_text
        r_top.font.size = Pt(12)
        r_top.font.bold = True
        r_top.font.color.rgb = pptx.dml.color.RGBColor(34, 197, 94)

        idx_box = slide.shapes.add_textbox(Pt(900), Pt(14), Pt(120), Pt(26))
        tf_idx = idx_box.text_frame
        tf_idx.clear()
        p_idx = tf_idx.paragraphs[0]
        r_idx = p_idx.add_run()
        r_idx.text = f'{slide_label} {idx:02d}' if slide_label else f'{idx:02d}'
        r_idx.font.size = Pt(11)
        r_idx.font.color.rgb = pptx.dml.color.RGBColor(168, 180, 209)

        title_box = slide.shapes.add_textbox(Pt(24), Pt(48), Pt(1240), Pt(76))
        tf_title = title_box.text_frame
        tf_title.clear()
        p_title = tf_title.paragraphs[0]
        r_title = p_title.add_run()
        r_title.text = title
        r_title.font.size = Pt(34)
        r_title.font.bold = True
        r_title.font.color.rgb = pptx.dml.color.RGBColor(233, 240, 255)

        body_box = slide.shapes.add_textbox(Pt(44), Pt(132), Pt(1180), Pt(520))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True
        tf_body.clear()

        bullet_lines = []
        for ln in body_lines:
            text_ln = ln.strip()
            if not text_ln:
                continue
            if text_ln.startswith('### '):
                bullet_lines.append(('sub', text_ln[4:].strip()))
            elif text_ln.startswith('- '):
                bullet_lines.append(('bullet', text_ln[2:].strip()))
            elif re.match(r'^\d+\.\s+', text_ln):
                bullet_lines.append(('bullet', re.sub(r'^\d+\.\s+', '', text_ln)))
            elif not text_ln.startswith('#'):
                bullet_lines.append(('bullet', text_ln))

        bullet_lines = bullet_lines[:10]

        first = True
        for kind, text_ln in bullet_lines:
            if first:
                p = tf_body.paragraphs[0]
                first = False
            else:
                p = tf_body.add_paragraph()

            run = p.add_run()
            if kind == 'sub':
                run.text = text_ln
                run.font.size = Pt(21)
                run.font.bold = True
                run.font.color.rgb = pptx.dml.color.RGBColor(142, 228, 255)
                p.level = 0
            else:
                run.text = f'• {text_ln}'
                run.font.size = Pt(20)
                run.font.color.rgb = pptx.dml.color.RGBColor(232, 237, 248)
                p.level = 0

        footer = slide.shapes.add_textbox(Pt(24), Pt(668), Pt(860), Pt(18))
        tf_footer = footer.text_frame
        tf_footer.clear()
        p_footer = tf_footer.paragraphs[0]
        r_footer = p_footer.add_run()
        r_footer.text = footer_text
        r_footer.font.size = Pt(10)
        r_footer.font.color.rgb = pptx.dml.color.RGBColor(120, 139, 170)

        gen = slide.shapes.add_textbox(Pt(900), Pt(668), Pt(360), Pt(18))
        tf_gen = gen.text_frame
        tf_gen.clear()
        p_gen = tf_gen.paragraphs[0]
        r_gen = p_gen.add_run()
        r_gen.text = generated_line
        r_gen.font.size = Pt(9)
        r_gen.font.color.rgb = pptx.dml.color.RGBColor(120, 139, 170)

    prs.save(str(pptx_path))
    return True


def main():
    args = parse_args()
    md_path = Path(args.input)
    base = Path(args.basename)

    html_path = base.with_name(base.name + '_presentation.html')
    pdf_path = base.with_suffix('.pdf')
    pptx_path = base.with_suffix('.pptx')

    text = md_path.read_text(encoding='utf-8')
    slides = parse_slides(text)

    generated_label, slide_label, timestamp = localized_defaults(
        args.lang,
        args.generated_label,
        args.slide_label,
    )
    generated_line = f"{generated_label}: {timestamp}"

    build_html(
        slides,
        html_path,
        lang=args.lang,
        badge_text=args.badge,
        slide_label=slide_label,
        generated_line=generated_line,
    )
    build_pdf(
        slides,
        pdf_path,
        slide_label=slide_label,
        generated_line=generated_line,
    )
    has_pptx = build_pptx(
        slides,
        pptx_path,
        args.footer,
        badge_text=args.badge,
        slide_label=slide_label,
        generated_line=generated_line,
    )

    print(f'HTML={html_path}')
    print(f'PDF={pdf_path}')
    if has_pptx:
        print(f'PPTX={pptx_path}')
    else:
        print('PPTX=skipped (python-pptx not installed)')


if __name__ == '__main__':
    main()
